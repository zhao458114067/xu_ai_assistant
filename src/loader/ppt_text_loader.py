from pptx import Presentation
from langchain.schema import Document
from langchain.document_loaders.base import BaseLoader
from typing import List


class PPTXTextLoader(BaseLoader):
    def __init__(self, file_path: str):
        self.file_path = file_path

    def load(self) -> List[Document]:
        prs = Presentation(self.file_path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        content = "\n".join(texts)
        return [Document(page_content=content, metadata={"source": self.file_path})]
